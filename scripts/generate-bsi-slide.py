#!/usr/bin/env python3
"""
Generate a PPTX slide for SP-047 (Secure Agentic AI Frameworks)
with BSI IT-Grundschutz control references instead of NIST 800-53.

All diagram elements are drawn as individual, editable PPTX shapes
(no embedded images). Every rectangle, label, badge, and arrow is
a separate shape you can select, move, recolour, or delete in
PowerPoint / Keynote / Google Slides.

SVG viewBox 960×720 → PPTX at SCALE = 0.01 in/px  (100 SVG px = 1 inch)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

WORKSPACE = Path(__file__).resolve().parent.parent
OUTPUT = WORKSPACE / "SP-047-BSI-Grundschutz.pptx"

# ── Coordinate mapping ──────────────────────────────────────────────
# SVG viewBox 960×720 on a 13.333×7.5 widescreen slide.
# SCALE chosen so diagram fills height well; centred horizontally.
SCALE  = 0.01        # 1 SVG px = 0.01 inches
X_OFF  = 1.87        # horizontal offset to centre 9.6" in 13.333"
Y_OFF  = 0.15        # vertical offset

def sx(x):        return Inches(x * SCALE + X_OFF)
def sy(y):        return Inches(y * SCALE + Y_OFF)
def sw(w):        return Inches(w * SCALE)
def sh(h):        return Inches(h * SCALE)

# ── OSA palette ─────────────────────────────────────────────────────
C_DARKEST  = RGBColor(0x00, 0x17, 0x1F)
C_DARK     = RGBColor(0x00, 0x34, 0x59)
C_MID      = RGBColor(0x00, 0x7E, 0xA7)
C_LIGHT    = RGBColor(0x00, 0xA8, 0xE8)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_BG       = RGBColor(0xF8, 0xFA, 0xFC)
C_RED      = RGBColor(0xC0, 0x39, 0x2B)
C_RED_DARK = RGBColor(0x7F, 0x1D, 0x1D)
C_GREY     = RGBColor(0x64, 0x74, 0x8B)
C_GREY_L   = RGBColor(0x94, 0xA3, 0xB8)
C_TEXT     = RGBColor(0x33, 0x41, 0x55)
C_TEXT_L   = RGBColor(0x47, 0x55, 0x69)
C_BORDER   = RGBColor(0xE2, 0xE8, 0xF0)

# ── Helpers ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, *, fill=None, border=None, border_w=1.5,
             dash=False, rx=0):
    """Add a rectangle shape.  rx is ignored visually (PPTX uses MSO_SHAPE)."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, sx(x), sy(y), sw(w), sh(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_w)
        if dash:
            s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, text, *, size=8, bold=False, color=C_TEXT,
             anchor="middle", w=200, h=20):
    """Add a text box.  anchor='middle' centres; 'start' left-aligns."""
    # SVG text x,y is baseline; offset up ~70 % of font size for top-left
    y_adj = y - size * 1.0
    tb = slide.shapes.add_textbox(sx(x - (w/2 if anchor == "middle" else 0)),
                                  sy(y_adj), sw(w), sh(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER if anchor == "middle" else PP_ALIGN.LEFT
    return tb


def add_badge(slide, x, y, label, *, fill_rgb=C_MID, text_rgb=C_MID,
              w=48, h=18, font_size=7):
    """BSI control badge (rounded pill)."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               sx(x), sy(y), sw(w), sh(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_rgb
    # Make fill semi-transparent (30 %)
    s.fill.fore_color.brightness = 0.7
    s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_rgb
    # Vertical centre
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return s


def add_arrow(slide, x1, y1, x2, y2, *, color=C_MID, width=1.5, dash=False):
    """Straight connector with arrowhead at end."""
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                   sx(x1), sy(y1), sx(x2), sy(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    # End arrow
    end = c.line._ln  # lxml element
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tail = etree.SubElement(end, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return c


def add_line(slide, x1, y1, x2, y2, *, color=C_MID, width=1.5, dash=False):
    """Straight line without arrowhead."""
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                   sx(x1), sy(y1), sx(x2), sy(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return c


# ═══════════════════════════════════════════════════════════════════
#  BUILD THE SLIDE
# ═══════════════════════════════════════════════════════════════════

def build(slide):

    # ── Background ──
    add_rect(slide, 0, 0, 960, 720, fill=C_BG, rx=8)

    # ═══ 1. GOVERNANCE BAR (top) ═══════════════════════════════════
    add_rect(slide, 16, 16, 928, 60, fill=C_DARKEST, rx=6)
    add_text(slide, 480, 38, "AGENT LIFECYCLE & DEPLOYMENT GOVERNANCE",
             size=10, bold=True, color=C_WHITE, w=600)
    # All original badges were non-critical → bar is clean, title only

    # ═══ 2. AGENTIC FRAMEWORK zone (left) ═════════════════════════
    add_rect(slide, 16, 92, 200, 310, border=C_MID, border_w=2, dash=True, rx=6)
    add_rect(slide, 16, 92, 200, 28, fill=C_MID, rx=6)
    add_text(slide, 116, 110, "AGENTIC FRAMEWORK",
             size=9, bold=True, color=C_WHITE, w=180)

    # Framework name chips
    for lbl, bx in [("LangChain", 30), ("CrewAI", 120)]:
        add_rect(slide, bx, 132, 86 if lbl == "LangChain" else 82, 20, fill=C_BG, rx=3)
        add_text(slide, bx + (43 if lbl == "LangChain" else 41), 146, lbl,
                 size=7, bold=True, color=C_MID, w=80)
    for lbl, bx in [("AutoGen", 30), ("LangGraph", 120)]:
        add_rect(slide, bx, 156, 86 if lbl == "AutoGen" else 82, 20, fill=C_BG, rx=3)
        add_text(slide, bx + (43 if lbl == "AutoGen" else 41), 170, lbl,
                 size=7, bold=True, color=C_MID, w=80)

    # Orchestrator box
    add_rect(slide, 48, 192, 40, 48, border=C_MID, border_w=2, rx=3)
    add_text(slide, 68, 254, "Orchestrator", size=8, bold=True, color=C_TEXT, w=80)

    # Sub-agent boxes
    add_rect(slide, 120, 196, 28, 32, border=C_MID, border_w=1.5, rx=2)
    add_rect(slide, 154, 196, 28, 32, border=C_MID, border_w=1.5, rx=2)
    add_rect(slide, 188, 196, 16, 32, border=C_MID, border_w=1, rx=2)
    add_text(slide, 154, 254, "Sub-Agents", size=8, bold=True, color=C_TEXT, w=80)

    # Delegation arrow
    add_arrow(slide, 92, 212, 116, 212, color=C_MID)

    # Context isolation boundary
    add_rect(slide, 30, 270, 172, 28, fill=C_BG, border=C_DARK, border_w=1, rx=4)
    add_text(slide, 116, 284, "Context Isolation Boundary",
             size=7.5, bold=True, color=C_DARK, w=170)
    add_text(slide, 116, 294, "Inter-agent messages \u2260 shared context",
             size=6, color=C_GREY, w=170)

    # Framework zone badges (BSI)
    add_badge(slide, 50, 306, "ORP.4", w=44, h=16, font_size=7)
    add_badge(slide, 100, 306, "OPS.1.1.5", w=56, h=16, font_size=6.5)

    # Descriptive text
    add_text(slide, 116, 348, "Each agent: distinct identity,",
             size=6.5, color=C_TEXT_L, w=160)
    add_text(slide, 116, 360, "scoped context, auditable chain.",
             size=6.5, color=C_TEXT_L, w=160)
    add_text(slide, 116, 372, "Proposer \u2260 Approver.",
             size=6.5, color=C_TEXT_L, w=160)

    # ═══ 3. EXECUTION SANDBOX (centre) ════════════════════════════
    add_rect(slide, 232, 92, 296, 310, border=C_DARK, border_w=2, dash=True, rx=6)
    add_rect(slide, 232, 92, 296, 28, fill=C_DARK, rx=6)
    add_text(slide, 380, 110, "EXECUTION SANDBOX",
             size=9, bold=True, color=C_WHITE, w=250)

    # Container icon (simplified)
    add_rect(slide, 260, 130, 24, 16, border=C_DARK, border_w=1.5, rx=2)
    add_text(slide, 272, 160, "Ephemeral", size=7, bold=True, color=C_TEXT, w=60)
    add_text(slide, 272, 170, "Container", size=7, bold=True, color=C_TEXT, w=60)

    # Process Isolation box
    add_rect(slide, 300, 130, 100, 48, fill=C_BG, border=C_DARK, border_w=1, rx=3)
    add_text(slide, 350, 146, "Process Isolation",
             size=7, bold=True, color=C_DARK, w=96)
    add_text(slide, 350, 158, "Read-only rootfs",
             size=6, color=C_GREY, w=96)
    add_text(slide, 350, 168, "Non-root \u00b7 Ephemeral \u00b7 Signed",
             size=6, color=C_GREY, w=96)

    # Resource Limits box
    add_rect(slide, 416, 130, 100, 48, fill=C_BG, border=C_DARK, border_w=1, rx=3)
    add_text(slide, 466, 146, "Resource Limits",
             size=7, bold=True, color=C_DARK, w=96)
    add_text(slide, 466, 158, "Token budget \u00b7 CPU/mem",
             size=6, color=C_GREY, w=96)
    add_text(slide, 466, 168, "Max time \u00b7 Cost ceiling",
             size=6, color=C_GREY, w=96)

    # Guardrails Engine bar
    add_rect(slide, 248, 192, 264, 36, fill=C_BG, border=C_DARK, border_w=1, rx=4)
    add_text(slide, 310, 208, "\U0001f6e1 Guardrails Engine",
             size=8, bold=True, color=C_DARK, anchor="start", w=160)
    add_text(slide, 286, 220,
             "Input filter \u00b7 Output filter \u00b7 Action policy \u00b7 Injection detection",
             size=6.5, color=C_GREY, anchor="start", w=250)
    add_badge(slide, 456, 198, "APP.3.1", w=48, h=16, font_size=7)

    # Tool Registry bar
    add_rect(slide, 248, 238, 264, 36, fill=C_BG, border=C_DARK, border_w=1, rx=4)
    add_text(slide, 310, 254, "\U0001f4cb Tool Registry",
             size=8, bold=True, color=C_DARK, anchor="start", w=160)
    add_text(slide, 286, 266,
             "Catalogued \u00b7 Risk-rated \u00b7 Versioned \u00b7 Per-agent authorisation",
             size=6.5, color=C_GREY, anchor="start", w=250)
    add_badge(slide, 456, 244, "ORP.4", w=48, h=16, font_size=7)

    # Network Policy bar
    add_rect(slide, 248, 284, 264, 36, fill=C_BG, border=C_DARK, border_w=1, rx=4)
    add_text(slide, 310, 300, "\U0001f512 Network Policy",
             size=8, bold=True, color=C_DARK, anchor="start", w=160)
    add_text(slide, 286, 312,
             "Egress whitelist \u00b7 No unauth outbound \u00b7 Infra-level enforcement",
             size=6.5, color=C_GREY, anchor="start", w=250)
    add_badge(slide, 456, 290, "NET.1.1", w=48, h=16, font_size=7)

    # Sandbox bottom badges
    add_badge(slide, 290, 330, "ORP.4", w=44, h=16, font_size=7)
    add_badge(slide, 340, 330, "NET.1.2", w=48, h=16, font_size=7)

    # Sandbox descriptive text
    add_text(slide, 380, 368,
             "Ephemeral containers. Signed images. Least privilege.",
             size=6.5, color=C_TEXT_L, w=270)
    add_text(slide, 380, 380,
             "Every tool invocation through the registry gate.",
             size=6.5, color=C_TEXT_L, w=270)
    add_text(slide, 380, 392,
             "Circuit breakers halt runaway agents.",
             size=6.5, color=C_TEXT_L, w=270)

    # ═══ 4. RAG PIPELINE (bottom-left) ════════════════════════════
    add_rect(slide, 16, 418, 200, 160, fill=C_WHITE, border=C_BORDER, border_w=1.5, rx=6)
    add_rect(slide, 16, 418, 200, 28, fill=C_DARK, rx=6)
    add_text(slide, 116, 436, "RAG PIPELINE",
             size=9, bold=True, color=C_WHITE, w=180)

    # Documents box
    add_rect(slide, 38, 458, 28, 32, border=C_DARK, border_w=1.5, rx=2)
    add_text(slide, 52, 504, "Documents", size=7, color=C_TEXT, w=50)

    # Arrow doc → vector
    add_arrow(slide, 70, 474, 90, 474, color=C_MID)

    # Vector Store (simplified as rect)
    add_rect(slide, 96, 458, 36, 32, border=C_DARK, border_w=1.5, rx=2)
    add_text(slide, 114, 470, "V", size=9, bold=True, color=C_DARK, w=20)
    add_text(slide, 114, 504, "Vector Store", size=7, color=C_TEXT, w=60)

    # Arrow vector → retrieval
    add_arrow(slide, 136, 474, 152, 474, color=C_MID)

    # Retrieval Filter
    add_rect(slide, 156, 458, 48, 32, fill=C_BG, border=C_DARK, border_w=1, rx=3)
    add_text(slide, 180, 472, "Retrieval", size=6, bold=True, color=C_DARK, w=44)
    add_text(slide, 180, 482, "Filter", size=6, bold=True, color=C_DARK, w=44)
    add_text(slide, 180, 504, "Authz Gate", size=7, color=C_TEXT, w=60)

    # RAG badges
    add_badge(slide, 42, 520, "ORP.4", w=44, h=16, font_size=7)
    add_badge(slide, 92, 520, "APP.3.1", w=48, h=16, font_size=7)

    add_text(slide, 116, 556, "Encrypted. Provenance-tracked.",
             size=6.5, color=C_TEXT_L, w=160)
    add_text(slide, 116, 568, "Role-filtered retrieval.",
             size=6.5, color=C_TEXT_L, w=160)

    # ═══ 5. ENTERPRISE TOOLS (right) ══════════════════════════════
    add_rect(slide, 544, 92, 200, 168, fill=C_WHITE, border=C_BORDER, border_w=1.5, rx=6)
    add_rect(slide, 544, 92, 200, 28, fill=C_DARK, rx=6)
    add_text(slide, 644, 110, "ENTERPRISE TOOLS",
             size=9, bold=True, color=C_WHITE, w=180)

    # Tool icons (simplified rects)
    add_rect(slide, 562, 134, 36, 28, border=C_DARK, border_w=1.5, rx=2)
    add_text(slide, 580, 150, "{ }", size=8, bold=True, color=C_DARK, w=30)
    add_text(slide, 580, 176, "APIs", size=7, color=C_TEXT, w=40)

    add_rect(slide, 630, 134, 28, 28, border=C_DARK, border_w=1.5, rx=2)
    add_text(slide, 644, 150, "DB", size=8, bold=True, color=C_DARK, w=24)
    add_text(slide, 644, 176, "Databases", size=7, color=C_TEXT, w=60)

    add_rect(slide, 692, 134, 36, 28, border=C_DARK, border_w=1.5, rx=2)
    add_text(slide, 710, 150, ">_", size=8, bold=True, color=C_DARK, w=30)
    add_text(slide, 710, 176, "Code Exec", size=7, color=C_TEXT, w=60)

    # MCP Servers bar
    add_rect(slide, 562, 190, 168, 28, fill=C_BG, border=C_MID, border_w=1, rx=4)
    add_text(slide, 646, 206, "MCP Servers \u00b7 Plugins \u00b7 Functions",
             size=7.5, bold=True, color=C_MID, w=160)
    add_text(slide, 646, 216, "Each vetted, versioned, risk-rated",
             size=6, color=C_GREY, w=160)

    # Enterprise tool badges
    add_badge(slide, 572, 228, "APP.3.1", w=48, h=16, font_size=7)
    add_badge(slide, 626, 228, "ORP.4", w=44, h=16, font_size=7)

    # ═══ 6. AI MODEL PROVIDER (right-bottom) ══════════════════════
    add_rect(slide, 544, 276, 200, 126, fill=C_WHITE, border=C_BORDER, border_w=1.5, rx=6)
    add_rect(slide, 544, 276, 200, 28, fill=C_DARK, rx=6)
    add_text(slide, 644, 294, "AI MODEL PROVIDER",
             size=9, bold=True, color=C_WHITE, w=180)

    # Cloud icon (simplified)
    add_rect(slide, 614, 324, 44, 28, border=C_DARK, border_w=1.5, rx=14)
    add_text(slide, 636, 354, "Model API", size=8, color=C_TEXT, w=60)
    add_text(slide, 636, 366, "(Anthropic / OpenAI / etc.)",
             size=7, color=C_GREY, w=140)

    # Provider badges
    add_badge(slide, 572, 376, "APP.3.1", w=48, h=16, font_size=7)
    add_badge(slide, 626, 376, "NET.1.1", w=48, h=16, font_size=7)

    # ═══ 7. THREAT LANDSCAPE (far right) ══════════════════════════
    add_rect(slide, 760, 92, 184, 310, fill=RGBColor(0xFF, 0xF5, 0xF5),
             border=C_RED, border_w=1.5, dash=True, rx=6)
    add_rect(slide, 760, 92, 184, 28, fill=C_RED, rx=6)
    add_text(slide, 852, 110, "THREAT LANDSCAPE",
             size=9, bold=True, color=C_WHITE, w=170)

    add_text(slide, 775, 140, "OWASP Agentic Top 10",
             size=7.5, bold=True, color=C_RED, anchor="start", w=160)

    threats = [
        (160, "ASI-01  Agent Goal Hijack"),
        (176, "ASI-02  Tool Misuse"),
        (192, "ASI-03  Identity & Privilege Abuse"),
        (208, "ASI-04  Cascading Hallucination"),
        (224, "ASI-05  Memory Corruption"),
        (240, "ASI-09  Excessive Autonomy"),
        (256, "ASI-10  Cross-Agent Trust"),
    ]
    for ty, txt in threats:
        add_text(slide, 775, ty, txt, size=7, color=C_RED_DARK, anchor="start", w=160)

    # Separator
    add_line(slide, 775, 270, 930, 270, color=C_RED, width=0.5, dash=False)

    add_text(slide, 775, 288, "Additional Threats",
             size=7.5, bold=True, color=C_RED, anchor="start", w=160)
    extra_threats = [
        (306, "RAG retrieval poisoning"),
        (322, "Framework supply chain CVEs"),
        (338, "Runaway resource consumption"),
        (354, "Guardrail bypass (indirect)"),
        (370, "Shadow agent deployment"),
        (386, "Agent credential theft"),
    ]
    for ty, txt in extra_threats:
        add_text(slide, 775, ty, txt, size=7, color=C_RED_DARK, anchor="start", w=160)

    # Threat arrows (red dashed, pointing left)
    for ty in [170, 250, 340]:
        add_line(slide, 760, ty, 748, ty, color=C_RED, width=1, dash=True)

    # ═══ 8. KILL SWITCH (centre-bottom) ═══════════════════════════
    add_rect(slide, 232, 418, 296, 60, fill=C_WHITE, border=C_RED, border_w=2, rx=6)
    add_rect(slide, 232, 418, 296, 28, fill=C_RED, rx=6)
    add_text(slide, 380, 436, "KILL SWITCH & INCIDENT RESPONSE",
             size=9, bold=True, color=C_WHITE, w=280)
    add_text(slide, 380, 462,
             "Halt agent \u00b7 Revoke credentials \u00b7 Preserve state \u00b7 Audit trail",
             size=7, color=C_RED_DARK, w=280)
    add_badge(slide, 355, 468, "DER.2.1", w=52, h=16, font_size=7,
              fill_rgb=C_RED, text_rgb=C_RED)

    # ═══ 9. MONITORING BAR (bottom) ═══════════════════════════════
    add_rect(slide, 16, 596, 928, 60, fill=C_DARKEST, rx=6)
    add_text(slide, 480, 618, "CONTINUOUS MONITORING, AUDIT & OBSERVABILITY",
             size=10, bold=True, color=C_WHITE, w=600)
    # Monitoring badges
    add_badge(slide, 180, 628, "OPS.1.1.5", w=60, h=18, font_size=7,
              fill_rgb=C_MID, text_rgb=C_LIGHT)
    add_badge(slide, 260, 628, "DER.1", w=52, h=18, font_size=7,
              fill_rgb=C_MID, text_rgb=C_LIGHT)
    add_badge(slide, 330, 628, "DER.2.1", w=52, h=18, font_size=7,
              fill_rgb=C_MID, text_rgb=C_LIGHT)
    add_text(slide, 580, 640,
             "Token usage \u00b7 Cost tracking \u00b7 Anomaly detection \u00b7 Delegation chains",
             size=7, color=C_GREY, w=350)

    # ═══ 10. DATA FLOW ARROWS ═════════════════════════════════════
    # Framework → Sandbox
    add_arrow(slide, 216, 210, 232, 210, color=C_MID)
    # Sandbox → Enterprise Tools (tool invocations)
    add_arrow(slide, 528, 165, 544, 165, color=C_MID)
    # Enterprise Tools → Sandbox (data return)
    add_arrow(slide, 544, 200, 528, 200, color=C_LIGHT, dash=True)
    # Sandbox → AI Provider
    add_arrow(slide, 528, 320, 544, 320, color=C_MID)
    # AI Provider → Sandbox (responses)
    add_arrow(slide, 544, 350, 528, 350, color=C_LIGHT, dash=True)
    # RAG → Sandbox (context)
    add_arrow(slide, 180, 418, 282, 402, color=C_LIGHT, dash=True)

    # ═══ 11. LEGEND ═══════════════════════════════════════════════
    add_text(slide, 16, 682, "SP-047: Secure Agentic AI Frameworks",
             size=9, bold=True, color=C_TEXT, anchor="start", w=300)
    add_text(slide, 16, 698,
             "12 BSI IT-Grundschutz controls (critical, agentic AI specific)  \u00b7  "
             "Authors: Aurelius, Vitruvius  \u00b7  Draft  \u00b7  2026-02-17",
             size=8, color=C_GREY, anchor="start", w=500)

    # Legend key: lines
    add_line(slide, 464, 690, 496, 690, color=C_DARK, width=1.5)
    add_text(slide, 502, 694, "Control flow", size=7.5, color=C_GREY, anchor="start", w=80)
    add_line(slide, 564, 690, 596, 690, color=C_LIGHT, width=1.5, dash=True)
    add_text(slide, 602, 694, "Data / response", size=7.5, color=C_GREY, anchor="start", w=80)
    add_badge(slide, 674, 682, "BSI.x", w=40, h=16, font_size=7)
    add_text(slide, 720, 694, "BSI Grundschutz Baustein",
             size=7.5, color=C_GREY, anchor="start", w=120)
    add_line(slide, 830, 690, 862, 690, color=C_RED, width=1, dash=True)
    add_text(slide, 868, 694, "Threat vector", size=7.5, color=C_GREY, anchor="start", w=80)

    add_text(slide, 944, 694, "opensecurityarchitecture.org",
             size=8, color=C_GREY_L, anchor="start", w=160)

    # ═══ 12. MAPPING TABLE — moved to slide 2 ═════════════════════


# ── NIST → BSI mapping data ─────────────────────────────────────
NIST_TO_BSI = {
    "AC-03": "ORP.4",     "AC-04": "ORP.4",
    "AC-05": "ORP.4",     "AC-06": "ORP.4",
    "AU-02": "OPS.1.1.5", "AU-03": "OPS.1.1.5",
    "CA-07": "DER.1",     "CM-07": "NET.1.2",
    "IR-04": "DER.2.1",   "SA-09": "APP.3.1",
    "SC-07": "NET.1.1",   "SI-10": "APP.3.1",
}

BSI_TITLES = {
    "APP.3.1":   "Web Apps & Web Services",
    "DER.1":     "Security Event Detection",
    "DER.2.1":   "Incident Management",
    "NET.1.1":   "Network Architecture & Design",
    "NET.1.2":   "Network Management",
    "OPS.1.1.5": "Logging",
    "ORP.4":     "Identity & Access Management",
}

BSI_DESCRIPTIONS = {
    "APP.3.1":   "Covers external service integration, input validation, and web application security controls for agent-to-tool and agent-to-model communication.",
    "DER.1":     "Continuous monitoring of agent behaviour, anomaly detection on token usage, API patterns, and resource consumption.",
    "DER.2.1":   "Agent-specific incident runbooks: kill switch activation, credential revocation, blast radius assessment, audit trail preservation.",
    "NET.1.1":   "Network boundary enforcement for agent execution sandboxes — egress whitelisting, infra-level guardrails that agents cannot bypass.",
    "NET.1.2":   "Least functionality for agent runtimes — strip unnecessary capabilities, restrict outbound network, disable shell access unless required.",
    "OPS.1.1.5": "Event logging and audit records for agent sessions — token usage, tool invocations, delegation chains, cost tracking per workflow.",
    "ORP.4":     "Agent identity, access enforcement, least privilege, separation of duties — each agent gets distinct credentials, scoped tool access, proposer ≠ approver.",
}


def build_mapping_slide(slide):
    """Slide 2: BSI IT-Grundschutz mapping reference table."""

    # ── Title ──
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "BSI IT-Grundschutz Control Mapping — SP-047 Agentic AI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    p.alignment = PP_ALIGN.LEFT

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Critical controls only — mapped from NIST 800-53 Rev 5 via OSA framework coverage data"
    p2.font.size = Pt(12)
    p2.font.color.rgb = C_MID
    p2.alignment = PP_ALIGN.LEFT

    # ── Build unique BSI rows ──
    unique_bsi = {}
    for nist_id, bsi_id in sorted(NIST_TO_BSI.items()):
        unique_bsi.setdefault(bsi_id, []).append(nist_id)

    rows = len(unique_bsi) + 1  # +1 header
    cols = 4
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(0.8), Inches(1.5),
        Inches(11.7), Inches(0.55 * rows),
    )
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(1.4)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(6.6)
    table.columns[3].width = Inches(1.5)

    def set_cell(cell, text, *, bold=False, size=10, color=None, bg=None, align=PP_ALIGN.LEFT):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        p.alignment = align
        cell.text_frame.word_wrap = True

    # Header
    hdr_bg = C_DARK
    hdr_fg = C_WHITE
    set_cell(table.cell(0, 0), "BSI Baustein", bold=True, size=10, color=hdr_fg, bg=hdr_bg)
    set_cell(table.cell(0, 1), "Title", bold=True, size=10, color=hdr_fg, bg=hdr_bg)
    set_cell(table.cell(0, 2), "Agentic AI Relevance", bold=True, size=10, color=hdr_fg, bg=hdr_bg)
    set_cell(table.cell(0, 3), "NIST 800-53", bold=True, size=10, color=hdr_fg, bg=hdr_bg)

    # Data rows
    for i, (bsi_id, nist_ids) in enumerate(sorted(unique_bsi.items()), start=1):
        title = BSI_TITLES.get(bsi_id, bsi_id)
        desc = BSI_DESCRIPTIONS.get(bsi_id, "")
        nist_str = ", ".join(nist_ids)
        row_bg = RGBColor(0xF0, 0xF8, 0xFF) if i % 2 == 0 else None

        set_cell(table.cell(i, 0), bsi_id, bold=True, size=10, color=C_MID, bg=row_bg)
        set_cell(table.cell(i, 1), title, size=9, bg=row_bg)
        set_cell(table.cell(i, 2), desc, size=8, color=C_TEXT_L, bg=row_bg)
        set_cell(table.cell(i, 3), nist_str, size=8, color=C_GREY, bg=row_bg)

    # ── Footer ──
    fb = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
    fp = fb.text_frame.paragraphs[0]
    fp.text = (
        "opensecurityarchitecture.org  |  "
        "Source: data/framework-coverage/bsi-grundschutz.json  |  "
        "Draft 2026-02-17"
    )
    fp.font.size = Pt(9)
    fp.font.color.rgb = C_GREY_L
    fp.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Architecture diagram ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = C_WHITE
    print("Drawing slide 1 (architecture diagram)...")
    build(slide1)
    print(f"  → {sum(1 for _ in slide1.shapes)} shapes")

    # ── Slide 2: BSI mapping table ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = C_WHITE
    print("Drawing slide 2 (BSI mapping table)...")
    build_mapping_slide(slide2)
    print(f"  → {sum(1 for _ in slide2.shapes)} shapes")

    print(f"Saving → {OUTPUT}")
    prs.save(str(OUTPUT))
    print("Done!")


if __name__ == "__main__":
    main()
