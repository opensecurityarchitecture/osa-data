#!/usr/bin/env python3
"""
OSA Pattern Portfolio — Executive Brief (SABSA Conceptual Layer)
3 slides: Cover, Heat Map, SP-047 Pattern Brief.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = Path(__file__).resolve().parent.parent / "OSA-Executive-Brief.pptx"

# ── Palette ──
C_DARK    = RGBColor(0x00, 0x34, 0x59)
C_MID     = RGBColor(0x00, 0x7E, 0xA7)
C_LIGHT   = RGBColor(0x00, 0xA8, 0xE8)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BG      = RGBColor(0xF8, 0xFA, 0xFC)
C_TEXT    = RGBColor(0x33, 0x41, 0x55)
C_GREY    = RGBColor(0x64, 0x74, 0x8B)
C_GREY_L  = RGBColor(0x94, 0xA3, 0xB8)
C_RED     = RGBColor(0xC0, 0x39, 0x2B)
C_AMBER   = RGBColor(0xE6, 0x7E, 0x22)
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)
C_DARKEST = RGBColor(0x00, 0x17, 0x1F)


def set_para(tf, text, *, size=10, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, idx=0):
    """Write text into paragraph idx of a text frame."""
    while len(tf.paragraphs) <= idx:
        tf.add_paragraph()
    p = tf.paragraphs[idx]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return p


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════

def build_cover(slide):
    # Full-width dark band
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.0),
                               Inches(13.333), Inches(3.4))
    s.fill.solid(); s.fill.fore_color.rgb = C_DARK
    s.line.fill.background(); s.shadow.inherit = False

    # Accent line
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.1),
                               Inches(0.08), Inches(1.6))
    a.fill.solid(); a.fill.fore_color.rgb = C_LIGHT
    a.line.fill.background(); a.shadow.inherit = False

    # Title
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    set_para(tf, "Open Security Architecture", size=16, color=C_LIGHT, idx=0)
    set_para(tf, "Pattern Portfolio", size=40, bold=True, color=C_WHITE, idx=1)
    set_para(tf, "Executive View", size=28, color=C_GREY_L, idx=2)

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.0), Inches(10), Inches(0.8))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    set_para(tf2, "SABSA Conceptual Layer — Business Risk & Capability Perspective",
             size=14, color=C_GREY_L)
    set_para(tf2, "48 security patterns mapped to business outcomes",
             size=12, color=C_GREY, idx=1)

    # Footer
    tb3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
    set_para(tb3.text_frame,
             "opensecurityarchitecture.org  |  Draft  |  February 2026",
             size=9, color=C_GREY_L, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 2 — PORTFOLIO HEAT MAP
# ═══════════════════════════════════════════════════════════════════

# Patterns grouped by business capability domain
DOMAINS = {
    "Identity & Access": [
        ("SP-010", "Identity Management"),
        ("SP-032", "Modern Authentication"),
        ("SP-033", "Passkey Authentication"),
        ("SP-037", "Privileged User Management"),
        ("SP-044", "SaaS Identity Lifecycle"),
    ],
    "Network & Perimeter": [
        ("SP-006", "Wireless Private Network"),
        ("SP-007", "Wireless Public Hotspot"),
        ("SP-008", "Public Web Server"),
        ("SP-016", "DMZ Module"),
        ("SP-017", "Secure Network Zone"),
        ("SP-029", "Zero Trust Architecture"),
        ("SP-046", "External Attack Surface Mgmt"),
    ],
    "Application & Development": [
        ("SP-004", "SOA Publication & Location"),
        ("SP-005", "SOA Internal Service Usage"),
        ("SP-012", "Secure SDLC"),
        ("SP-028", "Secure DevOps Pipeline"),
        ("SP-030", "API Security"),
        ("SP-041", "Secure App Baseline"),
    ],
    "Data Protection": [
        ("SP-003", "Privacy Mobile Device"),
        ("SP-013", "Data Security"),
        ("SP-019", "Secure File Exchange"),
        ("SP-020", "Email TLS"),
        ("SP-039", "Client-Side Encryption"),
        ("SP-040", "Post-Quantum Crypto"),
    ],
    "AI & Emerging Tech": [
        ("SP-027", "Secure AI Integration"),
        ("SP-045", "AI Governance"),
        ("SP-047", "Agentic AI Frameworks"),
    ],
    "Operations & Resilience": [
        ("SP-023", "Industrial Control Systems"),
        ("SP-025", "Advanced Monitoring"),
        ("SP-031", "Security Monitoring & Response"),
        ("SP-034", "Cyber Resilience"),
        ("SP-035", "Offensive Security Testing"),
        ("SP-036", "Incident Response"),
        ("SP-038", "Vulnerability Management"),
    ],
    "Governance & Compliance": [
        ("SP-014", "Awareness & Training"),
        ("SP-018", "ISMS"),
        ("SP-042", "Third Party Risk Mgmt"),
        ("SP-043", "Security Metrics"),
    ],
    "Endpoint & Workplace": [
        ("SP-001", "Client Module"),
        ("SP-002", "Server Module"),
        ("SP-009", "Generic Pattern"),
        ("SP-011", "Cloud Computing"),
        ("SP-015", "Secure Remote Working"),
        ("SP-021", "Realtime Collaboration"),
        ("SP-022", "Board Room"),
        ("SP-024", "iPhone"),
        ("SP-026", "PCI Full Environment"),
    ],
}


def build_heatmap(slide):
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    set_para(tb.text_frame, "Security Pattern Portfolio — Business Capability View",
             size=20, bold=True, color=C_DARK)

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.65), Inches(12), Inches(0.3))
    set_para(tb2.text_frame,
             "48 patterns across 8 capability domains  |  "
             "Each tile = one security pattern  |  "
             "Colour = business urgency",
             size=10, color=C_GREY)

    # Grid layout: 4 columns × 2 rows of domain blocks
    col_w = 3.05
    row_h = 2.7
    margin_x = 0.5
    margin_y = 1.15
    gap_x = 0.15
    gap_y = 0.15

    domain_list = list(DOMAINS.items())

    for di, (domain_name, patterns) in enumerate(domain_list):
        col = di % 4
        row = di // 4
        bx = margin_x + col * (col_w + gap_x)
        by = margin_y + row * (row_h + gap_y)

        # Domain header
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(bx), Inches(by),
                                     Inches(col_w), Inches(0.32))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = C_DARK
        hdr.line.fill.background(); hdr.shadow.inherit = False
        hdr.text_frame.word_wrap = False
        p = hdr.text_frame.paragraphs[0]
        p.text = domain_name
        p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER

        # Pattern tiles
        tile_cols = 3
        tw = (col_w - 0.1) / tile_cols
        th = 0.38

        for pi, (pid, pname) in enumerate(patterns):
            tc = pi % tile_cols
            tr = pi // tile_cols
            tx = bx + 0.05 + tc * tw
            ty = by + 0.38 + tr * (th + 0.04)

            # Urgency colour: AI domain = amber, default = mid blue
            if "AI" in domain_name:
                tile_fill = C_AMBER
                tile_text = C_WHITE
            elif domain_name == "Operations & Resilience":
                tile_fill = C_MID
                tile_text = C_WHITE
            else:
                tile_fill = RGBColor(0xE0, 0xF0, 0xFA)
                tile_text = C_DARK

            tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(tx), Inches(ty),
                                          Inches(tw - 0.04), Inches(th))
            tile.fill.solid(); tile.fill.fore_color.rgb = tile_fill
            tile.line.fill.background(); tile.shadow.inherit = False
            tf = tile.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(3); tf.margin_right = Pt(3)
            tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
            p0 = tf.paragraphs[0]
            p0.text = pid
            p0.font.size = Pt(6); p0.font.bold = True; p0.font.color.rgb = tile_text
            p0.alignment = PP_ALIGN.LEFT
            p1 = tf.add_paragraph()
            p1.text = pname
            p1.font.size = Pt(6); p1.font.color.rgb = tile_text
            p1.alignment = PP_ALIGN.LEFT

    # Legend
    ly = 6.85
    for lbl, clr, lx in [("Emerging / High urgency", C_AMBER, 0.5),
                          ("Active ops priority", C_MID, 3.5),
                          ("Foundational / Steady-state", RGBColor(0xE0, 0xF0, 0xFA), 6.5)]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(lx), Inches(ly), Inches(0.15), Inches(0.15))
        dot.fill.solid(); dot.fill.fore_color.rgb = clr
        dot.line.fill.background(); dot.shadow.inherit = False
        ltb = slide.shapes.add_textbox(Inches(lx + 0.22), Inches(ly - 0.02),
                                       Inches(2.5), Inches(0.2))
        set_para(ltb.text_frame, lbl, size=8, color=C_GREY)

    # Footer
    ftb = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
    set_para(ftb.text_frame,
             "opensecurityarchitecture.org  |  Colour reflects strategic urgency, not maturity  |  Draft Feb 2026",
             size=8, color=C_GREY_L, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 3 — PATTERN BRIEF: SP-047
# ═══════════════════════════════════════════════════════════════════

def build_brief_047(slide):
    # ── Header bar ──
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C_DARK
    hdr.line.fill.background(); hdr.shadow.inherit = False

    # Accent
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(0.15), Inches(0.06), Inches(0.8))
    acc.fill.solid(); acc.fill.fore_color.rgb = C_AMBER
    acc.line.fill.background(); acc.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(0.15), Inches(8), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    set_para(tf, "SP-047", size=12, color=C_LIGHT, idx=0)
    set_para(tf, "Secure Agentic AI Frameworks", size=24, bold=True, color=C_WHITE, idx=1)
    set_para(tf, "Pattern Brief — Executive View", size=11, color=C_GREY_L, idx=2)

    # Urgency badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(10.5), Inches(0.3), Inches(2.2), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = C_AMBER
    badge.line.fill.background(); badge.shadow.inherit = False
    btf = badge.text_frame; btf.word_wrap = False
    set_para(btf, "EMERGING — HIGH URGENCY", size=10, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    # ── 6 cards in 2×3 grid ──
    cards = [
        ("Business Risk",
         "What happens without this pattern?",
         C_RED,
         "Uncontrolled AI agents accessing enterprise tools, APIs, and data "
         "without security guardrails.\n\n"
         "- Agent goal hijack via prompt injection → unauthorised actions on production systems\n"
         "- Runaway agent costs → unbudgeted API spend (10–100× expected)\n"
         "- Data exfiltration through agent tool chains → regulatory exposure (GDPR, DORA)\n"
         "- Shadow agent deployments outside governance → unknown attack surface\n"
         "- Reputational damage from agent-generated external communications"),

        ("Capability Enabled",
         "What can the business do with this pattern?",
         C_GREEN,
         "Deploy agentic AI at enterprise scale with confidence.\n\n"
         "- Governed adoption of LangChain, CrewAI, AutoGen for automation\n"
         "- Multi-agent workflows for complex tasks (code review, compliance, analysis)\n"
         "- Safe connection of AI agents to enterprise APIs, databases, code execution\n"
         "- Clear autonomy boundaries: what agents decide vs. what humans approve\n"
         "- Auditable agent actions for regulatory and internal assurance"),

        ("Investment Profile",
         "What does adoption require?",
         C_AMBER,
         "Complexity: HIGH  |  Effort: 6–12 months to baseline\n\n"
         "- Container platform (K8s/ECS) for agent execution isolation\n"
         "- Centralised tool registry and MCP server governance\n"
         "- Guardrails tooling (NeMo Guardrails / custom policy engine)\n"
         "- Agent-aware monitoring and observability stack\n"
         "- Security team AI/ML competence development"),

        ("Dependencies",
         "What must be in place first?",
         C_MID,
         "Prerequisites:\n"
         "  SP-027  Secure AI Integration (individual agent controls)\n"
         "  SP-029  Zero Trust Architecture (network segmentation)\n\n"
         "Complementary:\n"
         "  SP-045  AI Governance (policy and ethics layer)\n"
         "  SP-030  API Security (tool endpoint protection)\n"
         "  SP-031  Security Monitoring (detection infrastructure)\n"
         "  SP-012  Secure SDLC (agent workflow CI/CD)"),

        ("Residual Risk",
         "What remains after adoption?",
         C_RED,
         "- Non-deterministic agent behaviour — testing reduces but cannot eliminate\n"
         "- Framework supply chain churn — LangChain et al. evolve faster than controls\n"
         "- Novel prompt injection vectors — adversarial research outpaces defences\n"
         "- Multi-agent emergent behaviour — agent interactions can produce unexpected outcomes\n"
         "- Regulatory uncertainty — DORA, EU AI Act interpretations still evolving"),

        ("Regulatory Relevance",
         "Which regulations does this address?",
         C_DARK,
         "- EU AI Act — high-risk AI system requirements, transparency, human oversight\n"
         "- DORA — ICT risk management for autonomous AI in financial services\n"
         "- GDPR — data processing by AI agents, automated decision-making (Art. 22)\n"
         "- FINMA guidance — operational risk controls for autonomous systems\n"
         "- NIST AI RMF — Agentic AI Profile (2025)\n"
         "- OWASP Top 10 for Agentic Applications (2025)"),
    ]

    card_w = 3.95
    card_h = 2.55
    gap_x = 0.16
    gap_y = 0.16
    start_x = 0.5
    start_y = 1.3

    for ci, (title, subtitle, accent_color, body) in enumerate(cards):
        col = ci % 3
        row = ci // 3
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + gap_y)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(cx), Inches(cy),
                                      Inches(card_w), Inches(card_h))
        card.fill.solid(); card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        card.line.width = Pt(1)
        card.shadow.inherit = False

        # Accent stripe at top
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(cx), Inches(cy),
                                        Inches(card_w), Inches(0.06))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = accent_color
        stripe.line.fill.background(); stripe.shadow.inherit = False

        # Title
        ttb = slide.shapes.add_textbox(Inches(cx + 0.15), Inches(cy + 0.12),
                                       Inches(card_w - 0.3), Inches(0.55))
        ttf = ttb.text_frame; ttf.word_wrap = True
        set_para(ttf, title, size=11, bold=True, color=C_DARK, idx=0)
        set_para(ttf, subtitle, size=8, color=C_GREY, idx=1)

        # Body
        btb = slide.shapes.add_textbox(Inches(cx + 0.15), Inches(cy + 0.65),
                                       Inches(card_w - 0.3), Inches(card_h - 0.75))
        btf = btb.text_frame; btf.word_wrap = True
        # Split body into lines for better control
        for li, line in enumerate(body.split('\n')):
            set_para(btf, line, size=8, color=C_TEXT, idx=li)

    # Footer
    ftb = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
    set_para(ftb.text_frame,
             "opensecurityarchitecture.org  |  SABSA Conceptual Layer  |  Draft Feb 2026",
             size=8, color=C_GREY_L, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for name, builder in [("Cover", build_cover),
                          ("Heat Map", build_heatmap),
                          ("SP-047 Brief", build_brief_047)]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = C_BG
        print(f"  Slide: {name}...")
        builder(slide)

    prs.save(str(OUTPUT))
    print(f"Done → {OUTPUT}")


if __name__ == "__main__":
    main()
