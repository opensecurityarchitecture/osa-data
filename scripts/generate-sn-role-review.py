#!/usr/bin/env python3
"""Generate a 5-slide PPTX summarizing ServiceNow role concept review (German)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Palette
DARK_NAVY = RGBColor(0x00, 0x17, 0x1F)
NAVY = RGBColor(0x00, 0x34, 0x59)
TEAL = RGBColor(0x00, 0x7E, 0xA7)
BLUE = RGBColor(0x00, 0xA8, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x27, 0xAE, 0x60)
MID_GRAY = RGBColor(0x7F, 0x8C, 0x8D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, radius=Inches(0.1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_footer(slide, text="ServiceNow Rollen- & Berechtigungskonzept \u2014 Review & Empfehlungen"):
    add_rect(slide, Inches(0), H - Inches(0.45), W, Inches(0.45), NAVY)
    add_text_box(slide, Inches(0.5), H - Inches(0.4), Inches(8), Inches(0.35), text, font_size=9, color=TEAL)


def add_slide_number(slide, num):
    add_text_box(slide, W - Inches(1), H - Inches(0.4), Inches(0.6), Inches(0.35), str(num), font_size=9, color=TEAL, alignment=PP_ALIGN.RIGHT)


# ── Slide 1: Title ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

# Accent bar
add_rect(slide, Inches(0), Inches(2.8), W, Inches(0.06), BLUE)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "ServiceNow", font_size=20, color=TEAL, bold=False)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Rollen- & Berechtigungskonzept", font_size=44, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.8),
             "Review, Sicherheitsprinzipien & Empfehlungen", font_size=22, color=BLUE, bold=False)
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             "Februar 2026  |  Vertraulich", font_size=14, color=MID_GRAY)

# Key findings box
box = add_shape(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(1.4), NAVY)
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Zusammenfassung: Das bestehende Konzept bietet eine gute operative Basis mit ~30 Rollensets,"
p.font.size = Pt(13)
p.font.color.rgb = LIGHT_GRAY
p.font.name = "Calibri"
add_para(tf, "Umgebungs-spezifischer Provisionierung und automatisiertem Lifecycle-Management.", font_size=13, color=LIGHT_GRAY)
add_para(tf, "Es fehlen jedoch zentrale Sicherheitsprinzipien, Aufgabentrennung (SoD) und eine nachhaltige Governance.", font_size=13, color=BLUE, bold=True)

add_footer(slide)
add_slide_number(slide, 1)


# ── Slide 2: Sicherheitsprinzipien ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Fehlende Sicherheitsprinzipien", font_size=32, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), BLUE)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Das Konzept definiert keine expliziten Sicherheitsprinzipien als Steuerungsgrundlage.", font_size=14, color=MID_GRAY)

principles = [
    ("Least Privilege", "Nur minimale Rechte pro Rolle", "Teilweise \u2014 einige Rollen\nsind sehr breit gefasst", AMBER),
    ("Aufgabentrennung\n(SoD)", "Konfligierende Rechte in\ngetrennten Rollen", "Nicht adressiert \u2014\nkeine SoD-Matrix", RED),
    ("Need-to-Know", "Datenzugriff an Funktion\ngekoppelt", "Nur implizit \u2014 keine\nDatenklassifizierung", AMBER),
    ("Defense in Depth", "ACLs + Rollen + Row-Level\nals Schichten", "ACLs kaum\ndokumentiert", RED),
    ("Fail Secure", "Default Deny; Zugriff\nmuss explizit erteilt werden", "Nicht definiert", RED),
    ("Accountability", "Jede Aktion r\u00fcckverfolgbar\nauf eine Person", "Shared Accounts\nnicht geregelt", AMBER),
]

col_w = Inches(1.85)
gap = Inches(0.17)
start_x = Inches(0.8)
start_y = Inches(1.9)
card_h = Inches(4.8)

for i, (title, desc, status, status_color) in enumerate(principles):
    x = start_x + i * (col_w + gap)

    card = add_shape(slide, x, start_y, col_w, card_h, NAVY)
    card.text_frame.word_wrap = True

    # Title
    add_text_box(slide, x + Inches(0.15), start_y + Inches(0.15), col_w - Inches(0.3), Inches(0.8),
                 title, font_size=14, color=BLUE, bold=True)

    # Description
    add_text_box(slide, x + Inches(0.15), start_y + Inches(1.0), col_w - Inches(0.3), Inches(1.2),
                 desc, font_size=11, color=LIGHT_GRAY)

    # Status indicator
    status_bar = add_rect(slide, x, start_y + Inches(2.5), col_w, Inches(0.04), status_color)

    # Status text
    add_text_box(slide, x + Inches(0.15), start_y + Inches(2.7), col_w - Inches(0.3), Inches(1.0),
                 status, font_size=11, color=status_color, bold=True)

    # Bottom label
    label = "Kritisch" if status_color == RED else "Verbesserungsbedarf"
    add_text_box(slide, x + Inches(0.15), start_y + card_h - Inches(0.6), col_w - Inches(0.3), Inches(0.4),
                 label, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer(slide)
add_slide_number(slide, 2)


# ── Slide 3: Kritische Befunde ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Kritische Befunde", font_size=32, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), RED)

findings = [
    (
        "Admin-Rolle als Fallback",
        "In den Men\u00fc-\u00c4nderungen werden systematisch OOTB-Rollen entfernt und durch \u00abadmin\u00bb ersetzt "
        "(SLA, Problem, Change, CMDB, IntegrationHub, etc.). Die admin-Rolle umgeht alle ACLs \u2014 "
        "jeder Benutzer mit dieser Rolle kann alles sehen und \u00e4ndern.",
        "OOTB-Rollen beibehalten oder durch massgeschneiderte, eingeschr\u00e4nkte Rollen ersetzen \u2014 niemals admin als Ersatz.",
        RED,
    ),
    (
        "Process Manager Mega-Rolle",
        "Die Rolle \u00abProcess Manager Asset & SCM\u00bb aggregiert ~16 technische Rollen inkl. itil_admin, "
        "workflow_admin, fm_billing_admin und cmdb_inst_admin. Ein Benutzer kann damit Incidents l\u00f6schen, "
        "Workflows \u00e4ndern UND Finanzdaten einsehen.",
        "Aufteilen in mindestens 3 Unter-Rollen: CMDB Operations, CMDB Administration, Financial/Service Ownership.",
        RED,
    ),
    (
        "Keine Aufgabentrennung (SoD)",
        "Es gibt keine Dokumentation, welche Rollenkombinationen Konflikte erzeugen. "
        "Z.B. Change-Antragsteller vs. Change-Genehmiger, Impersonator + operative Rolle, "
        "Dispatcher + Agent, Procurement + Asset Owner.",
        "SoD-Matrix erstellen und mit Excluded Roles oder Business Rules auf sys_user_has_role durchsetzen.",
        RED,
    ),
]

card_y = Inches(1.4)
card_h_each = Inches(1.8)
card_gap = Inches(0.15)

for i, (title, desc, recommendation, color) in enumerate(findings):
    y = card_y + i * (card_h_each + card_gap)

    # Left accent bar
    add_rect(slide, Inches(0.8), y, Inches(0.06), card_h_each, color)

    # Card background
    card = add_shape(slide, Inches(0.95), y, Inches(11.5), card_h_each, NAVY)

    # Title
    add_text_box(slide, Inches(1.15), y + Inches(0.12), Inches(11), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)

    # Description
    add_text_box(slide, Inches(1.15), y + Inches(0.55), Inches(11), Inches(0.65),
                 desc, font_size=11, color=LIGHT_GRAY)

    # Recommendation
    rec_box = add_text_box(slide, Inches(1.15), y + Inches(1.25), Inches(11), Inches(0.45),
                           f"Empfehlung: {recommendation}", font_size=11, color=GREEN, bold=False)

add_footer(slide)
add_slide_number(slide, 3)


# ── Slide 4: Governance & Lifecycle ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Governance & Lifecycle-L\u00fccken", font_size=32, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), BLUE)

# Left column: Current state
left_x = Inches(0.8)
right_x = Inches(7.0)
col_width = Inches(5.8)

# Current state card
add_text_box(slide, left_x, Inches(1.3), col_width, Inches(0.4),
             "Ist-Zustand (vorhanden)", font_size=16, color=GREEN, bold=True)

current_items = [
    "Automatisierte Birthright-Rollen via IAM/SSO + MFA",
    "Access Groups f\u00fcr Rollenzuweisung (gruppenbasiert)",
    "Umgebungs-spezifische Provisionierung (DEV/QUAL/PROD)",
    "Inactive User Cleanup mit Decision Table & Benachrichtigung",
    "ACL-Testbenutzer pro Rolleset auf DEV",
    "Men\u00fc-/Modul-Sichtbarkeitsmatrix dokumentiert",
]

y_pos = Inches(1.8)
for item in current_items:
    box = add_shape(slide, left_x, y_pos, col_width, Inches(0.42), NAVY)
    add_text_box(slide, left_x + Inches(0.15), y_pos + Inches(0.05), col_width - Inches(0.3), Inches(0.35),
                 f"\u2713  {item}", font_size=11, color=GREEN)
    y_pos += Inches(0.5)

# Missing state card
add_text_box(slide, right_x, Inches(1.3), col_width, Inches(0.4),
             "Soll-Zustand (fehlend)", font_size=16, color=RED, bold=True)

missing_items = [
    "Genehmigungs-Workflow: Wer genehmigt welche Rollenvergabe?",
    "Rollen-Rezertifizierung: Quartalsweise f\u00fcr privilegierte Rollen",
    "Break-Glass-Verfahren: Notfallzugriff ausserhalb der Gesch\u00e4ftszeiten",
    "Audit & Monitoring: \u00dcberwachung von sys_user_has_role",
    "Rollen-Eigent\u00fcmer: Verantwortlicher pro Rollen-Definition",
    "Vollst\u00e4ndigkeit: TBD-Eintr\u00e4ge, leere Rollen (snc_external, Mobile Admin)",
]

y_pos = Inches(1.8)
for item in missing_items:
    box = add_shape(slide, right_x, y_pos, col_width, Inches(0.42), NAVY)
    add_text_box(slide, right_x + Inches(0.15), y_pos + Inches(0.05), col_width - Inches(0.3), Inches(0.35),
                 f"\u2717  {item}", font_size=11, color=RED)
    y_pos += Inches(0.5)

# Bottom recommendation
rec_box = add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5), NAVY)
tf = rec_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Empfehlung: Rollen-Lifecycle definieren"
p.font.size = Pt(14)
p.font.color.rgb = BLUE
p.font.bold = True
p.font.name = "Calibri"
add_para(tf, "Pro Rollenset dokumentieren: Eigent\u00fcmer (wer \u00e4ndert die Definition), Genehmiger (wer erteilt Zugriff),", font_size=11, color=LIGHT_GRAY)
add_para(tf, "Rezertifizierungs-Kadenz (Q f\u00fcr privilegiert, J f\u00fcr Standard), maximale Zuweisungsdauer (v.a. admin, impersonator).", font_size=11, color=LIGHT_GRAY)
add_para(tf, "Inactive User Cleanup um Rollen-Rezertifizierung f\u00fcr aktive Benutzer erg\u00e4nzen.", font_size=11, color=LIGHT_GRAY)

add_footer(slide)
add_slide_number(slide, 4)


# ── Slide 5: Priorisierte Massnahmen ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Priorisierte Massnahmen", font_size=32, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), BLUE)

actions = [
    ("P1", "Kritisch", RED, [
        "Admin-Fallback in Men\u00fc-Sichtbarkeit durch eingeschr\u00e4nkte Custom Roles ersetzen",
        "Process Manager Mega-Rolle in 3 Sub-Rollen aufteilen (Operations / Admin / Finance)",
        "SoD-Matrix erstellen und Enforcement via Excluded Roles oder Business Rules implementieren",
    ]),
    ("P2", "Hoch", AMBER, [
        "Sicherheitsprinzipien-Sektion im Dokument verankern (Least Privilege, SoD, Need-to-Know, etc.)",
        "Rollen-Lifecycle dokumentieren: Antrag, Genehmigung, Rezertifizierung, Break-Glass",
        "Alle TBD-Eintr\u00e4ge und leeren Rollendefinitionen aufl\u00f6sen",
    ]),
    ("P3", "Mittel", TEAL, [
        "Rollen-Hierarchie / Containment-Diagramm erstellen",
        "ACL-Strategie und Custom-ACL-Inventar dokumentieren",
        "Tippfehler korrigieren und Sprache vereinheitlichen (DE oder EN)",
    ]),
    ("P4", "Niedrig", MID_GRAY, [
        "Rollen-Rezertifizierung f\u00fcr aktive Benutzer einf\u00fchren (erg\u00e4nzend zum Inactive Cleanup)",
    ]),
]

y_start = Inches(1.4)
for i, (prio, label, color, items) in enumerate(actions):
    y = y_start + i * Inches(1.45)
    row_h = Inches(1.3)

    # Priority badge
    badge = add_shape(slide, Inches(0.8), y, Inches(0.9), Inches(0.45), color)
    badge.text_frame.paragraphs[0].text = prio
    badge.text_frame.paragraphs[0].font.size = Pt(14)
    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.name = "Calibri"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.text_frame.word_wrap = False

    # Label
    add_text_box(slide, Inches(1.85), y, Inches(1.2), Inches(0.45),
                 label, font_size=13, color=color, bold=True)

    # Action items
    for j, item in enumerate(items):
        bullet_y = y + Inches(0.5) + j * Inches(0.28)
        add_text_box(slide, Inches(1.85), bullet_y, Inches(10.5), Inches(0.28),
                     f"\u2022  {item}", font_size=11, color=LIGHT_GRAY)

add_footer(slide)
add_slide_number(slide, 5)


# ── Slide 6: ABAC in ServiceNow ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "ABAC \u2014 Attribute-Based Access Control", font_size=32, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), BLUE)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "ServiceNow nutzt ein hybrides RBAC+ABAC-Modell \u2014 ihr setzt ABAC bereits implizit ein.", font_size=14, color=MID_GRAY)

# ── Left: How ACLs work (3-tier) ──
left_x = Inches(0.8)
left_w = Inches(5.8)

add_text_box(slide, left_x, Inches(1.8), left_w, Inches(0.4),
             "ACL-Pr\u00fcfung: 3 Ebenen (AND-verkn\u00fcpft)", font_size=16, color=BLUE, bold=True)

acl_tiers = [
    ("1. Rolle (RBAC)", "Benutzer hat die erforderliche Rolle\nz.B. itil, asset, cmdb_read", TEAL),
    ("2. Condition (ABAC deklarativ)", "Attribut-Pr\u00fcfung ohne Script\nz.B. assignment_group IS one of my groups", BLUE),
    ("3. Script (ABAC programmatisch)", "GlideScript-basierte Logik\nz.B. current.managed_by == gs.getUserID()", RGBColor(0x8E, 0x44, 0xAD)),
]

tier_y = Inches(2.3)
for label, desc, color in acl_tiers:
    # Colored left bar
    add_rect(slide, left_x, tier_y, Inches(0.06), Inches(1.0), color)
    # Card
    card = add_shape(slide, left_x + Inches(0.15), tier_y, left_w - Inches(0.15), Inches(1.0), NAVY)
    add_text_box(slide, left_x + Inches(0.3), tier_y + Inches(0.08), left_w - Inches(0.5), Inches(0.35),
                 label, font_size=13, color=color, bold=True)
    add_text_box(slide, left_x + Inches(0.3), tier_y + Inches(0.42), left_w - Inches(0.5), Inches(0.55),
                 desc, font_size=11, color=LIGHT_GRAY)
    tier_y += Inches(1.1)

# ── Right: ABAC mechanisms ──
right_x = Inches(7.0)
right_w = Inches(5.8)

add_text_box(slide, right_x, Inches(1.8), right_w, Inches(0.4),
             "ABAC-Mechanismen in ServiceNow", font_size=16, color=BLUE, bold=True)

mechanisms = [
    ("Before Query Business Rules", "Row-Level Security \u2014 filtert Datens\u00e4tze vor Anzeige"),
    ("Domain Separation", "Zugriff basierend auf Firmen-/Dom\u00e4nen-Zugeh\u00f6rigkeit"),
    ("Data Classification", "Felder mit Sensitivity-Label versehen (seit Washington DC)"),
    ("Contextual Security", "Zugriff abh\u00e4ngig von Netzwerk, IP, MFA-Status"),
    ("Delegated ACLs", "Zugriff basierend auf Beziehungen (Manager von, Mitglied von)"),
]

mech_y = Inches(2.3)
for mech_title, mech_desc in mechanisms:
    card = add_shape(slide, right_x, mech_y, right_w, Inches(0.52), NAVY)
    add_text_box(slide, right_x + Inches(0.15), mech_y + Inches(0.03), right_w - Inches(0.3), Inches(0.25),
                 mech_title, font_size=12, color=TEAL, bold=True)
    add_text_box(slide, right_x + Inches(0.15), mech_y + Inches(0.27), right_w - Inches(0.3), Inches(0.25),
                 mech_desc, font_size=10, color=LIGHT_GRAY)
    mech_y += Inches(0.58)

# ── Bottom: What you already do + recommendation ──
bottom_y = Inches(5.5)

# "Already using" box
already_box = add_shape(slide, Inches(0.8), bottom_y, Inches(5.8), Inches(1.2), NAVY)
tf = already_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Bereits im Einsatz (implizit)"
p.font.size = Pt(13)
p.font.color.rgb = GREEN
p.font.bold = True
p.font.name = "Calibri"
add_para(tf, "\"Can update assets and CIs if: Managed by Group,", font_size=11, color=LIGHT_GRAY)
add_para(tf, "Managed by, Owned by\" \u2014 das ist ABAC via ACL-Conditions", font_size=11, color=LIGHT_GRAY)
add_para(tf, "oder Before Query Business Rules.", font_size=11, color=LIGHT_GRAY)

# Recommendation box
rec_box = add_shape(slide, Inches(7.0), bottom_y, Inches(5.8), Inches(1.2), NAVY)
tf = rec_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Empfehlungen"
p.font.size = Pt(13)
p.font.color.rgb = BLUE
p.font.bold = True
p.font.name = "Calibri"
add_para(tf, "\u2022  Im Konzept als \u00abhybrides RBAC+ABAC-Modell\u00bb deklarieren", font_size=11, color=LIGHT_GRAY)
add_para(tf, "\u2022  Attribut-Regeln pro Rolle inventarisieren (ACL + Business Rule)", font_size=11, color=LIGHT_GRAY)
add_para(tf, "\u2022  Data Classification f\u00fcr sensitive Felder einf\u00fchren", font_size=11, color=LIGHT_GRAY)

add_footer(slide)
add_slide_number(slide, 6)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/tobias.christen/Downloads/ServiceNow-Rollenkonzept-Review.pptx"
prs.save(out)
print(f"Saved → {out}")
