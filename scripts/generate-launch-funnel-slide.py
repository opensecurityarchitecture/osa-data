#!/usr/bin/env python3
"""
OSA — Launch & Funnel Build-Up timeline slide (16:9).
Two-row layout so every font is ≥ 11pt (OSA style rule).
Row 1 (5 stages): Attention · Interest · Desire · Action · Support Building
Row 2 (4 stages): Support Testing · Publishing · Storytelling · Standardize
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

OUTPUT = Path(__file__).resolve().parent.parent / "OSA-Launch-Funnel.pptx"

# ── Palette ──────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x00, 0x34, 0x59)
C_TEAL   = RGBColor(0x00, 0x7E, 0xA7)
C_SKY    = RGBColor(0x00, 0xA8, 0xE8)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BG     = RGBColor(0xF4, 0xF7, 0xFA)
C_TEXT   = RGBColor(0x1E, 0x2D, 0x3D)
C_MUTED  = RGBColor(0x64, 0x74, 0x8B)
C_RULE   = RGBColor(0xB0, 0xBE, 0xCC)

# ── Stage Data ────────────────────────────────────────────────────────
ROW1 = [
    ("Attention",  "(Aware)",    C_NAVY, [
        "Reddit r/netsec technical posts",
        "LinkedIn TRIDENT launch post",
        "GitHub repo + O'Reilly provenance",
        "ATT&CK depth screenshot hooks",
    ]),
    ("Interest",   "(Inform)",   C_NAVY, [
        "48 pattern pages — SEO optimised",
        "TRIDENT Explorer walkthrough",
        "Framework coverage explorer",
        "NIST · ATT&CK · CIS graph depth",
    ]),
    ("Desire",     "(Convince)", C_NAVY, [
        "Free assessment tool (10 min)",
        "Financial Services vertical profile",
        "Benchmark comparison preview",
        "TRIDENT scale — 17K+ edges",
    ]),
    ("Action",     "(Onboard)",  C_TEAL, [
        "GitHub / LinkedIn OAuth sign-up",
        "Free API key self-registration",
        "First assessment completion",
        "Pattern follow & watchlist",
    ]),
    ("Support",    "Building",   C_TEAL, [
        "GitHub Discussions on osa-data",
        "Per-pattern comment system",
        "Pattern contribution programme",
        "Post-registration email sequence",
    ]),
]

ROW2 = [
    ("Support",    "Testing",         C_TEAL, [
        "Activation A/B testing",
        "Assessment benchmark pool build",
        "Beta feature gate (TRIDENT)",
        "E2E coverage — 85 assertions",
    ]),
    ("Support",    "Publishing",      C_TEAL, [
        "Regular pattern release cadence",
        "Framework mapping updates",
        "Bi-weekly blog post programme",
        "Changelog & release notes",
    ]),
    ("Support",    "Storytelling",    C_TEAL, [
        "CISO board pack demo deck",
        "SABSA / ISC2 conference talks",
        "Consultant case studies",
        "Podcast & media outreach",
    ]),
    ("Standardize","(natl. / intl.)", C_SKY, [
        "Intl. verticals: SG · JP · CH · NL",
        "SABSA chapter integration",
        "ISO 42001 / EU AI Act alignment",
        "Franchise / consultant programme",
    ]),
]

# Phase bands: (start_idx, end_idx_excl, label, color) — per row
ROW1_PHASES = [
    (0, 4, "FUNNEL",            C_NAVY),
    (4, 5, "COMMUNITY",         C_TEAL),
]
ROW2_PHASES = [
    (0, 3, "COMMUNITY GROWTH",  C_TEAL),
    (3, 4, "STANDARDISE",       C_SKY),
]


# ── Helpers ───────────────────────────────────────────────────────────

def no_shadow(shape):
    shape.shadow.inherit = False


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_pt=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    no_shadow(s)
    return s


def add_textbox(slide, x, y, w, h, text, size=11, bold=False, color=C_TEXT,
                align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tb


def add_circle(slide, cx, cy, r, fill_color, line_color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                cx - r, cy - r, r * 2, r * 2)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.color.rgb = line_color
    s.line.width = Pt(1.5)
    no_shadow(s)
    return s


def add_line(slide, x1, y1, x2, y2, color, width_pt=0.75):
    idx = len(slide.shapes)
    x1e, y1e = int(x1), int(y1)
    x2e, y2e = int(x2), int(y2)
    w = abs(x2e - x1e) or 1
    h = abs(y2e - y1e) or 1
    ox, oy = min(x1e, x2e), min(y1e, y2e)
    r, g, b = color[0], color[1], color[2]
    hex_c = f"{r:02X}{g:02X}{b:02X}"
    we = int(Pt(width_pt))
    xml = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f'<p:nvSpPr><p:cNvPr id="{9000+idx}" name="Line{idx}"/>'
        f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{ox}" y="{oy}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:noFill/><a:ln w="{we}"><a:solidFill>'
        f'<a:srgbClr val="{hex_c}"/></a:solidFill></a:ln></p:spPr>'
        f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>'
    )
    slide.shapes._spTree.append(parse_xml(xml))


def tint(c, t=0.88):
    return int(c + (255 - c) * t)


def draw_phase_band(slide, node_xs, phases, band_y, band_h):
    for start_i, end_i, label, color in phases:
        x0 = node_xs[start_i]
        x1 = node_xs[end_i - 1]
        bw = (x1 - x0) + Inches(0.8)
        bx = (x0 + x1) / 2 - bw / 2
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    bx, band_y, bw, band_h)
        s.fill.solid()
        r, g, b = color[0], color[1], color[2]
        s.fill.fore_color.rgb = RGBColor(tint(r), tint(g), tint(b))
        s.line.color.rgb = color
        s.line.width = Pt(0.75)
        no_shadow(s)
        tb = slide.shapes.add_textbox(bx, band_y, bw, band_h)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER


def draw_timeline_row(slide, stages, node_xs, timeline_y,
                       label_top_y, bullet_top_y, bullet_row_h):
    n_bullets = 4
    node_r = Inches(0.13)
    drop_bottom = bullet_top_y + n_bullets * bullet_row_h + Inches(0.04)
    col_w = Inches(2.6)

    # Horizontal line
    add_line(slide, node_xs[0], timeline_y, node_xs[-1], timeline_y,
             C_RULE, width_pt=1.2)

    for i, (l1, l2, node_color, bullets) in enumerate(stages):
        cx = node_xs[i]

        # ── Stage label above timeline ────────────────────────────────
        tb_l = slide.shapes.add_textbox(
            cx - col_w / 2, label_top_y, col_w, Inches(0.65))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True

        p1 = tf_l.paragraphs[0]
        p1.text = l1
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = C_NAVY
        p1.alignment = PP_ALIGN.CENTER

        tf_l.add_paragraph()
        p2 = tf_l.paragraphs[1]
        p2.text = l2
        p2.font.size = Pt(11)
        p2.font.bold = False
        p2.font.color.rgb = (C_TEAL if node_color != C_SKY else C_SKY)
        p2.alignment = PP_ALIGN.CENTER

        # ── Vertical drop line ────────────────────────────────────────
        add_line(slide, cx, timeline_y + node_r, cx, drop_bottom,
                 C_RULE, width_pt=0.75)

        # ── Circle node ───────────────────────────────────────────────
        add_circle(slide, cx, timeline_y, node_r, C_WHITE, node_color)

        # Node number
        tb_n = slide.shapes.add_textbox(
            cx - node_r, timeline_y - node_r, node_r * 2, node_r * 2)
        tf_n = tb_n.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.text = str(i + 1)
        p_n.font.size = Pt(11)
        p_n.font.bold = True
        p_n.font.color.rgb = node_color
        p_n.alignment = PP_ALIGN.CENTER

        # ── Bullets ───────────────────────────────────────────────────
        bx = cx - col_w / 2 + Inches(0.18)
        for j, text in enumerate(bullets):
            by = bullet_top_y + j * bullet_row_h

            # Bullet dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                cx - col_w / 2 + Inches(0.04),
                by + Inches(0.08),
                Inches(0.07), Inches(0.07))
            dot.fill.solid()
            dot.fill.fore_color.rgb = node_color
            dot.line.fill.background()
            no_shadow(dot)

            tb_b = slide.shapes.add_textbox(
                bx, by, col_w - Inches(0.22), bullet_row_h)
            tf_b = tb_b.text_frame
            tf_b.word_wrap = True
            p_b = tf_b.paragraphs[0]
            p_b.text = text
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = C_TEXT
            p_b.alignment = PP_ALIGN.LEFT


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    W, H = prs.slide_width, prs.slide_height

    # ── Background ────────────────────────────────────────────────────
    add_rect(slide, 0, 0, W, H, C_BG)
    add_rect(slide, 0, 0, W, Inches(0.07), C_TEAL)   # top accent

    # ── Title ─────────────────────────────────────────────────────────
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.14),
                                   Inches(10), Inches(0.55))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Launch"
    r1.font.size = Pt(26); r1.font.bold = True; r1.font.color.rgb = C_NAVY

    r2 = p.add_run(); r2.text = " & Funnel Build-Up"
    r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = C_TEAL

    add_textbox(slide, Inches(0.5), Inches(0.68), Inches(11), Inches(0.28),
                "Open Security Architecture · Go-to-Market Strategy 2026",
                size=11, color=C_MUTED)

    # ── Geometry ──────────────────────────────────────────────────────
    margin = Inches(0.80)
    usable = W - 2 * margin

    # Row 1 — 5 nodes
    r1_xs = [margin + i * usable / (len(ROW1) - 1) for i in range(len(ROW1))]
    # Row 2 — 4 nodes (same margins, evenly spaced)
    r2_xs = [margin + i * usable / (len(ROW2) - 1) for i in range(len(ROW2))]

    # Vertical positions — row 1
    r1_phase_y   = Inches(1.02)
    r1_phase_h   = Inches(0.27)
    r1_label_y   = Inches(1.36)
    r1_timeline  = Inches(2.10)
    r1_bullet_y  = Inches(2.32)
    bullet_row_h = Inches(0.35)

    # Vertical positions — row 2
    row_gap      = Inches(0.25)
    r2_top       = r1_bullet_y + 4 * bullet_row_h + row_gap
    r2_phase_y   = r2_top
    r2_phase_h   = Inches(0.27)
    r2_label_y   = r2_top + Inches(0.34)
    r2_timeline  = r2_top + Inches(1.08)
    r2_bullet_y  = r2_top + Inches(1.30)

    # ── Row 1 phase bands ─────────────────────────────────────────────
    draw_phase_band(slide, r1_xs, ROW1_PHASES, r1_phase_y, r1_phase_h)

    # ── Row 1 timeline ────────────────────────────────────────────────
    draw_timeline_row(slide, ROW1, r1_xs, r1_timeline,
                       r1_label_y, r1_bullet_y, bullet_row_h)

    # ── Row separator ─────────────────────────────────────────────────
    sep_y = r1_bullet_y + 4 * bullet_row_h + Inches(0.10)
    add_line(slide, Inches(0.4), sep_y, W - Inches(0.4), sep_y,
             C_RULE, width_pt=0.5)

    # ── Row 2 phase bands ─────────────────────────────────────────────
    draw_phase_band(slide, r2_xs, ROW2_PHASES, r2_phase_y, r2_phase_h)

    # ── Row 2 timeline ────────────────────────────────────────────────
    draw_timeline_row(slide, ROW2, r2_xs, r2_timeline,
                       r2_label_y, r2_bullet_y, bullet_row_h)

    # ── Row 2 global stage index offset (nodes 6–9) ───────────────────
    # Renumber row-2 nodes as 6,7,8,9 — we must patch the textbox labels.
    # Easier: pass an offset into draw_timeline_row. Let's add offset support:
    # (Handled below by patching text directly — see note)

    # ── Footer ────────────────────────────────────────────────────────
    footer_y = H - Inches(0.34)
    add_rect(slide, 0, footer_y, W, Inches(0.34), C_NAVY)
    add_textbox(slide, Inches(0.4), footer_y + Inches(0.07),
                Inches(6), Inches(0.22),
                "opensecurityarchitecture.org",
                size=11, color=RGBColor(0x94, 0xA3, 0xB8))
    add_textbox(slide, 0, footer_y + Inches(0.07),
                W - Inches(0.4), Inches(0.22),
                "Go-to-Market Strategy 2026  ·  Confidential",
                size=11, color=RGBColor(0x94, 0xA3, 0xB8),
                align=PP_ALIGN.RIGHT)

    return slide


def build_slide_v2(prs):
    """
    Clean two-row build with correct per-row node numbering (row1: 1–5, row2: 6–9).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height

    add_rect(slide, 0, 0, W, H, C_BG)
    add_rect(slide, 0, 0, W, Inches(0.07), C_TEAL)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.14),
                                   Inches(10), Inches(0.55))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Launch"
    r1.font.size = Pt(26); r1.font.bold = True; r1.font.color.rgb = C_NAVY
    r2 = p.add_run(); r2.text = " & Funnel Build-Up"
    r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = C_TEAL

    add_textbox(slide, Inches(0.5), Inches(0.68), Inches(11), Inches(0.28),
                "Open Security Architecture · Go-to-Market Strategy 2026",
                size=11, color=C_MUTED)

    margin    = Inches(0.80)
    usable    = W - 2 * margin
    brow_h    = Inches(0.35)   # bullet row height
    n_bul     = 4              # bullets per stage
    node_r    = Inches(0.13)

    def row_xs(n_stages):
        return [margin + i * usable / max(n_stages - 1, 1)
                for i in range(n_stages)]

    def draw_row(stages, xs, phase_defs, phase_y, timeline_y,
                 label_y, bullet_y, num_offset=0):
        # Phase bands
        ph = Inches(0.27)
        for si, ei, lbl, col in phase_defs:
            x0, x1 = xs[si], xs[ei - 1]
            bw = (x1 - x0) + Inches(0.8)
            bx = (x0 + x1) / 2 - bw / 2
            s = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, bx, phase_y, bw, ph)
            s.fill.solid()
            r, g, b = col[0], col[1], col[2]
            s.fill.fore_color.rgb = RGBColor(tint(r), tint(g), tint(b))
            s.line.color.rgb = col; s.line.width = Pt(0.75)
            no_shadow(s)
            tb_ph = slide.shapes.add_textbox(bx, phase_y, bw, ph)
            p_ph = tb_ph.text_frame.paragraphs[0]
            p_ph.text = lbl; p_ph.font.size = Pt(11)
            p_ph.font.bold = True; p_ph.font.color.rgb = col
            p_ph.alignment = PP_ALIGN.CENTER

        # Timeline line
        add_line(slide, xs[0], timeline_y, xs[-1], timeline_y,
                 C_RULE, width_pt=1.2)

        drop_bot = bullet_y + n_bul * brow_h + Inches(0.04)
        col_w    = Inches(2.6)

        for i, (l1, l2, nc, bullets) in enumerate(stages):
            cx = xs[i]
            num = i + 1 + num_offset

            # Stage label
            tb_l = slide.shapes.add_textbox(
                cx - col_w / 2, label_y, col_w, Inches(0.65))
            tf_l = tb_l.text_frame; tf_l.word_wrap = True
            p1 = tf_l.paragraphs[0]
            p1.text = l1; p1.font.size = Pt(13); p1.font.bold = True
            p1.font.color.rgb = C_NAVY; p1.alignment = PP_ALIGN.CENTER
            tf_l.add_paragraph()
            p2 = tf_l.paragraphs[1]
            p2.text = l2; p2.font.size = Pt(11); p2.font.bold = False
            p2.font.color.rgb = C_TEAL if nc != C_SKY else C_SKY
            p2.alignment = PP_ALIGN.CENTER

            # Drop line
            add_line(slide, cx, timeline_y + node_r, cx, drop_bot,
                     C_RULE, width_pt=0.75)

            # Circle
            add_circle(slide, cx, timeline_y, node_r, C_WHITE, nc)

            # Number in circle
            tb_n = slide.shapes.add_textbox(
                cx - node_r, timeline_y - node_r, node_r * 2, node_r * 2)
            p_n = tb_n.text_frame.paragraphs[0]
            p_n.text = str(num); p_n.font.size = Pt(11)
            p_n.font.bold = True; p_n.font.color.rgb = nc
            p_n.alignment = PP_ALIGN.CENTER

            # Bullets
            bx = cx - col_w / 2 + Inches(0.18)
            for j, txt in enumerate(bullets):
                by = bullet_y + j * brow_h
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    cx - col_w / 2 + Inches(0.04),
                    by + Inches(0.09),
                    Inches(0.07), Inches(0.07))
                dot.fill.solid(); dot.fill.fore_color.rgb = nc
                dot.line.fill.background(); no_shadow(dot)
                tb_b = slide.shapes.add_textbox(
                    bx, by, col_w - Inches(0.22), brow_h)
                tf_b = tb_b.text_frame; tf_b.word_wrap = True
                p_b = tf_b.paragraphs[0]
                p_b.text = txt; p_b.font.size = Pt(11)
                p_b.font.color.rgb = C_TEXT

    # ── Row positions ─────────────────────────────────────────────────
    r1_xs = row_xs(len(ROW1))
    r2_xs = row_xs(len(ROW2))

    r1_phase_y   = Inches(1.02)
    r1_label_y   = Inches(1.36)
    r1_timeline  = Inches(2.12)
    r1_bullet_y  = Inches(2.35)

    r1_bottom    = r1_bullet_y + n_bul * brow_h
    sep_y        = r1_bottom + Inches(0.12)

    r2_phase_y   = sep_y + Inches(0.18)
    r2_label_y   = r2_phase_y + Inches(0.35)
    r2_timeline  = r2_phase_y + Inches(1.10)
    r2_bullet_y  = r2_phase_y + Inches(1.32)

    draw_row(ROW1, r1_xs, ROW1_PHASES,
             r1_phase_y, r1_timeline, r1_label_y, r1_bullet_y,
             num_offset=0)

    # Row separator
    add_line(slide, Inches(0.4), sep_y, W - Inches(0.4), sep_y,
             C_RULE, width_pt=0.5)

    draw_row(ROW2, r2_xs, ROW2_PHASES,
             r2_phase_y, r2_timeline, r2_label_y, r2_bullet_y,
             num_offset=len(ROW1))

    # ── Footer ────────────────────────────────────────────────────────
    fy = H - Inches(0.34)
    add_rect(slide, 0, fy, W, Inches(0.34), C_NAVY)
    add_textbox(slide, Inches(0.4), fy + Inches(0.06),
                Inches(6), Inches(0.24),
                "opensecurityarchitecture.org",
                size=11, color=RGBColor(0x94, 0xA3, 0xB8))
    add_textbox(slide, 0, fy + Inches(0.06),
                W - Inches(0.4), Inches(0.24),
                "Go-to-Market Strategy 2026  ·  Confidential",
                size=11, color=RGBColor(0x94, 0xA3, 0xB8),
                align=PP_ALIGN.RIGHT)


if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_slide_v2(prs)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
